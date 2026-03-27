"""
API helpers and shared utilities for church_management.
"""
import datetime
import hashlib
import io
import os
import re
import secrets
from decimal import Decimal, InvalidOperation

from django.conf import settings
from django.db import transaction
from django.db.models import Q
from django.core.files.base import ContentFile
from django.core.mail import EmailMessage
from django.utils import timezone

from rest_framework.response import Response


def _is_admin_user(user):
    """Check if user has admin privileges."""
    if not user or not getattr(user, 'is_authenticated', False):
        return False
    if getattr(user, 'is_superuser', False):
        return True
    if getattr(user, 'is_staff', False):
        return True
    return getattr(user, 'role', None) in {'super_admin', 'pastor', 'admin', 'administrator'}


def _admin_recipients_qs():
    """Get queryset of admin users for notifications."""
    from .models import User
    return User.objects.filter(
        Q(is_superuser=True)
        | Q(is_staff=True)
        | Q(role__in=['super_admin', 'pastor', 'admin', 'administrator'])
    ).distinct()


def _notify_admins(title, message):
    """Send notification to all admins."""
    from .models import Notification
    for u in _admin_recipients_qs():
        Notification.objects.create(title=title, message=message, recipient=u)


def _notify_user(user, title, message):
    """Send notification to a specific user."""
    from .models import Notification
    if user and getattr(user, 'is_authenticated', False):
        Notification.objects.create(title=title, message=message, recipient=user)


def _client_ip(request):
    """Extract client IP from request."""
    xff = request.META.get('HTTP_X_FORWARDED_FOR')
    if xff:
        return xff.split(',')[0].strip()
    return request.META.get('REMOTE_ADDR')


def _safe_payload(data):
    """Sanitize payload for audit logging."""
    if data is None:
        return None
    if isinstance(data, dict):
        out = {}
        for k, v in data.items():
            if isinstance(k, str) and any(s in k.lower() for s in ('password', 'token', 'secret', 'key', 'auth')):
                out[k] = '***'
            else:
                out[k] = v
        return out
    return data


def _ensure_event_share_slug(event):
    """Generate unique share slug for public events."""
    if getattr(event, 'share_slug', None):
        return event.share_slug

    from .models import Event
    for _ in range(6):
        slug = secrets.token_urlsafe(18)
        if not Event.objects.filter(share_slug=slug).exists():
            event.share_slug = slug
            event.save(update_fields=['share_slug'])
            return slug

    slug = secrets.token_urlsafe(24)
    event.share_slug = slug
    event.save(update_fields=['share_slug'])
    return slug


def _create_approval_request(request, model, action, payload=None, target_object_id=None, object_repr=None):
    """Create approval request for non-admin actions."""
    from .models import ApprovalRequest
    ar = ApprovalRequest.objects.create(
        requested_by=request.user,
        model=model,
        action=action,
        payload=payload,
        target_object_id=target_object_id,
        object_repr=object_repr,
        ip_address=_client_ip(request),
    )
    _notify_admins(
        'Nouvelle demande d\'approbation',
        f"{request.user.username} demande l'approbation pour {action} {model}"
    )
    return ar


def _apply_approval_request(ar, request):
    """Apply approved action."""
    from .models import (
        Event, EventAttendanceAggregate, EventVisitorAggregate, EventLogisticsConsumption,
        EvangelismActivity, TrainingEvent, BaptismEvent, BaptismCandidate,
        FinancialTransaction, LogisticsItem, Attendance
    )

    model_map = {
        'Event': Event,
        'EventAttendanceAggregate': EventAttendanceAggregate,
        'EventVisitorAggregate': EventVisitorAggregate,
        'EventLogisticsConsumption': EventLogisticsConsumption,
        'EvangelismActivity': EvangelismActivity,
        'TrainingEvent': TrainingEvent,
        'BaptismEvent': BaptismEvent,
        'BaptismCandidate': BaptismCandidate,
        'FinancialTransaction': FinancialTransaction,
        'LogisticsItem': LogisticsItem,
        'Attendance': Attendance,
    }

    model_class = model_map.get(ar.model)
    if not model_class:
        raise ValueError(f"Unknown model: {ar.model}")

    payload = ar.payload or {}
    action = ar.action

    if action == 'create':
        instance = model_class.objects.create(**payload)
        ar.target_object_id = str(instance.pk)
    elif action == 'update' and ar.target_object_id:
        instance = model_class.objects.filter(id=ar.target_object_id).first()
        if instance:
            for k, v in payload.items():
                setattr(instance, k, v)
            instance.save()
    elif action == 'delete' and ar.target_object_id:
        model_class.objects.filter(id=ar.target_object_id).delete()


def _create_report_certificate(report_type, payload, user=None):
    """Create certificate for generated reports."""
    from .models import ReportCertificate
    cert = ReportCertificate.objects.create(
        report_type=report_type,
        payload=payload,
        created_by=user,
    )
    return cert


def _draw_authenticity_qr(c, verify_url, code, x, width, margin, qr_size_mm=26, qr_y_mm=34, text_top_mm=46):
    """Draw QR code for document verification."""
    import qrcode
    qr_img = qrcode.make(verify_url)
    try:
        qr_img = qr_img.convert('RGB')
    except Exception:
        pass

    from reportlab.lib.utils import ImageReader
    qx = width - margin - qr_size_mm * 3.7795275591
    qy = qr_y_mm * 3.7795275591

    c.setFillColorRGB(1, 1, 1)
    c.roundRect(qx, qy, qr_size_mm * 3.7795275591, qr_size_mm * 3.7795275591, 8, fill=1, stroke=0)
    c.drawImage(
        ImageReader(qr_img),
        qx,
        qy,
        qr_size_mm * 3.7795275591,
        qr_size_mm * 3.7795275591,
        preserveAspectRatio=True,
        mask='auto',
    )

    c.setFillColorRGB(0.07, 0.10, 0.16)
    c.setFont('Helvetica-Bold', 9)
    c.drawString(x, text_top_mm * 3.7795275591, f"Authenticité: {code}")


# PPTX helpers
def _pptx_add_title(prs, title_text, subtitle_text=None):
    """Add title slide to presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = title_text or ''
    if subtitle_text:
        subtitle.text = subtitle_text


def _pptx_add_announcement_slide(prs, header, text, number=None):
    """Add announcement slide to presentation."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except Exception:
        return

    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    header_text = header or 'Annonce'
    if number:
        header_text = f"{header_text} #{number}"
    title_shape.text = header_text

    tf = body_shape.text_frame
    tf.text = text or ''
