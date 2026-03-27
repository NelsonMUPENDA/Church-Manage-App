"""
Content API: Announcements, AnnouncementDecks, Documents, ChurchBiography, ChurchConsistory.
"""
import io
import re

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None
    Inches = None
    Pt = None

from django.db import transaction
from django.http import HttpResponse, FileResponse
from django.utils import timezone
from django.core.files.base import ContentFile

from rest_framework import viewsets
from rest_framework.decorators import action
from rest_framework.permissions import IsAuthenticated
from rest_framework.response import Response
from rest_framework.parsers import JSONParser, FormParser, MultiPartParser

from .helpers import _client_ip, _safe_payload
from ..models import (
    Announcement, AnnouncementComment, AnnouncementCommentLike, AnnouncementLike,
    AnnouncementDeck, AnnouncementDeckItem,
    Document, ChurchBiography, ChurchConsistory,
    AuditLogEntry
)
from ..permissions import IsAdminOrSuperAdmin, PublicReadAdminWrite
from ..serializers import (
    AnnouncementSerializer, AnnouncementCommentSerializer,
    AnnouncementDeckSerializer, AnnouncementDeckItemSerializer,
    DocumentSerializer, ChurchBiographySerializer, ChurchConsistorySerializer
)


class AnnouncementViewSet(viewsets.ModelViewSet):
    queryset = Announcement.objects.all().order_by('-published_date')
    serializer_class = AnnouncementSerializer
    permission_classes = [PublicReadAdminWrite]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def get_permissions(self):
        if getattr(self, 'action', None) in {'like', 'comments', 'comment_like'}:
            return [IsAuthenticated()]
        return super().get_permissions()

    def perform_create(self, serializer):
        serializer.save(author=self.request.user)

    @action(detail=True, methods=['post'], url_path='like')
    def like(self, request, pk=None):
        ann = self.get_object()
        existing = AnnouncementLike.objects.filter(announcement=ann, user=request.user).first()
        if existing:
            existing.delete()
            liked = False
        else:
            AnnouncementLike.objects.create(announcement=ann, user=request.user)
            liked = True
        return Response({'liked': liked, 'like_count': AnnouncementLike.objects.filter(announcement=ann).count()})

    @action(detail=True, methods=['get', 'post'], url_path='comments')
    def comments(self, request, pk=None):
        ann = self.get_object()
        if request.method == 'GET':
            qs = AnnouncementComment.objects.filter(announcement=ann).order_by('-created_at')
            return Response(AnnouncementCommentSerializer(qs, many=True, context={'request': request}).data)

        serializer = AnnouncementCommentSerializer(data=request.data, context={'request': request})
        serializer.is_valid(raise_exception=True)
        serializer.save(announcement=ann, author=request.user)
        return Response(serializer.data, status=201)

    @action(detail=True, methods=['post'], url_path=r'comments/(?P<comment_id>\d+)/like')
    def comment_like(self, request, pk=None, comment_id=None):
        ann = self.get_object()
        comment = AnnouncementComment.objects.filter(id=comment_id, announcement=ann).first()
        if not comment:
            return Response({'detail': 'commentaire introuvable'}, status=404)

        existing = AnnouncementCommentLike.objects.filter(comment=comment, user=request.user).first()
        if existing:
            existing.delete()
            liked = False
        else:
            AnnouncementCommentLike.objects.create(comment=comment, user=request.user)
            liked = True

        return Response({'liked': liked, 'like_count': AnnouncementCommentLike.objects.filter(comment=comment).count()})


class AnnouncementDeckViewSet(viewsets.ModelViewSet):
    queryset = AnnouncementDeck.objects.select_related('event', 'created_by').prefetch_related('items').all().order_by('-created_at', '-id')
    serializer_class = AnnouncementDeckSerializer
    permission_classes = [IsAdminOrSuperAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)

    def perform_update(self, serializer):
        serializer.save()

    @action(detail=True, methods=['get', 'post'], url_path='items')
    def items(self, request, pk=None):
        deck = self.get_object()

        if request.method == 'GET':
            qs = getattr(deck, 'items', None).all() if getattr(deck, 'items', None) is not None else AnnouncementDeckItem.objects.filter(deck=deck)
            return Response(AnnouncementDeckItemSerializer(qs, many=True).data)

        serializer = AnnouncementDeckItemSerializer(data=request.data)
        serializer.is_valid(raise_exception=True)
        serializer.save(deck=deck)
        return Response(serializer.data, status=201)

    @action(detail=True, methods=['post'], url_path='set-items')
    def set_items(self, request, pk=None):
        deck = self.get_object()
        raw = request.data.get('items')
        if not isinstance(raw, list):
            return Response({'detail': 'items doit être une liste.'}, status=400)

        cleaned = []
        for idx, it in enumerate(raw, start=1):
            if not isinstance(it, dict):
                continue
            text = (it.get('text') or '').strip()
            if not text:
                continue
            try:
                order = int(it.get('order') or idx)
            except Exception:
                order = idx
            cleaned.append({'order': max(1, order), 'text': text})

        with transaction.atomic():
            AnnouncementDeckItem.objects.filter(deck=deck).delete()
            for it in cleaned:
                AnnouncementDeckItem.objects.create(deck=deck, order=it['order'], text=it['text'])

        qs = AnnouncementDeckItem.objects.filter(deck=deck).order_by('order', 'id')
        return Response(AnnouncementDeckItemSerializer(qs, many=True).data)

    @action(detail=True, methods=['post'], url_path='generate')
    def generate(self, request, pk=None):
        deck = self.get_object()

        if Presentation is None:
            return Response({'detail': 'Génération PPTX indisponible: installe python-pptx sur le backend.'}, status=503)

        items = list(getattr(deck, 'items', None).all()) if getattr(deck, 'items', None) is not None else list(AnnouncementDeckItem.objects.filter(deck=deck))
        header = (getattr(deck, 'header_text', None) or '').strip()
        theme = (getattr(deck, 'theme_text', None) or '').strip()

        ev = getattr(deck, 'event', None)
        ev_title = getattr(ev, 'title', None) if ev else None
        ev_date = getattr(ev, 'date', None) if ev else None
        ev_time = getattr(ev, 'time', None) if ev else None
        ev_location = getattr(ev, 'location', None) if ev else None

        subtitle_parts = []
        if ev_date:
            subtitle_parts.append(str(ev_date))
        if ev_time:
            subtitle_parts.append(str(ev_time)[:5])
        if ev_location:
            subtitle_parts.append(str(ev_location))
        if theme:
            subtitle_parts.append(str(theme))
        subtitle = ' • '.join([p for p in subtitle_parts if str(p).strip()])

        prs = Presentation()
        title = (getattr(deck, 'title', None) or ev_title or 'Annonces')

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle

        # Header slide
        if header:
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = 'En-tête'
            slide.placeholders[1].text = header

        # Item slides
        for idx, it in enumerate(items, start=1):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = f"{title} #{idx}"
            slide.placeholders[1].text = getattr(it, 'text', '')

        buf = io.BytesIO()
        prs.save(buf)
        pptx = buf.getvalue()
        buf.close()

        stamp = timezone.now().strftime('%Y%m%d_%H%M%S')
        safe_title = re.sub(r'[^a-zA-Z0-9_-]+', '_', str(title))[:60].strip('_') or 'annonces'
        filename = f"deck_{safe_title}_{stamp}.pptx"

        deck.pptx_file.save(filename, ContentFile(pptx), save=False)
        deck.generated_at = timezone.now()
        deck.save(update_fields=['pptx_file', 'generated_at', 'updated_at'])

        resp = HttpResponse(pptx, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
        resp['Content-Disposition'] = f'attachment; filename="{filename}"'
        return resp

    @action(detail=True, methods=['get'], url_path='download')
    def download(self, request, pk=None):
        deck = self.get_object()
        f = getattr(deck, 'pptx_file', None)
        if not f:
            return Response({'detail': 'Fichier PPTX indisponible. Génère le deck.'}, status=404)
        try:
            return FileResponse(f.open('rb'), as_attachment=True, filename=getattr(f, 'name', None) or 'annonces.pptx')
        except Exception:
            return Response({'detail': 'Impossible de télécharger le fichier.'}, status=500)


class DocumentViewSet(viewsets.ModelViewSet):
    queryset = Document.objects.all().order_by('-uploaded_at', '-id')
    serializer_class = DocumentSerializer
    permission_classes = [IsAdminOrSuperAdmin]
    parser_classes = [JSONParser, FormParser, MultiPartParser]

    def perform_create(self, serializer):
        serializer.save(uploaded_by=self.request.user)


class ChurchBiographyViewSet(viewsets.ModelViewSet):
    queryset = ChurchBiography.objects.all()
    serializer_class = ChurchBiographySerializer
    permission_classes = [PublicReadAdminWrite]

    def get_queryset(self):
        qs = super().get_queryset()
        return qs.filter(is_active=True).order_by('-updated_at')

    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)
        self._log_action('create', serializer.instance)

    def perform_update(self, serializer):
        serializer.save()
        self._log_action('update', serializer.instance)

    def destroy(self, request, *args, **kwargs):
        obj = self.get_object()
        obj.is_active = False
        obj.save()
        self._log_action('delete', obj)
        return Response(status=204)

    def _log_action(self, action, obj):
        AuditLogEntry.objects.create(
            actor=getattr(self.request, 'user', None),
            action=action,
            model='ChurchBiography',
            object_id=str(getattr(obj, 'pk', '')),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'pk', '')),
            ip_address=_client_ip(self.request),
            payload={
                'title': getattr(obj, 'title', None),
                'is_active': getattr(obj, 'is_active', None),
            },
        )


class ChurchConsistoryViewSet(viewsets.ModelViewSet):
    queryset = ChurchConsistory.objects.all()
    serializer_class = ChurchConsistorySerializer
    permission_classes = [PublicReadAdminWrite]

    def get_queryset(self):
        qs = super().get_queryset()
        return qs.filter(is_active=True).order_by('-updated_at')

    def perform_create(self, serializer):
        serializer.save(created_by=self.request.user)
        self._log_action('create', serializer.instance)

    def perform_update(self, serializer):
        serializer.save()
        self._log_action('update', serializer.instance)

    def destroy(self, request, *args, **kwargs):
        obj = self.get_object()
        obj.is_active = False
        obj.save()
        self._log_action('delete', obj)
        return Response(status=204)

    def _log_action(self, action, obj):
        AuditLogEntry.objects.create(
            actor=getattr(self.request, 'user', None),
            action=action,
            model='ChurchConsistory',
            object_id=str(getattr(obj, 'pk', '')),
            object_repr=getattr(obj, 'title', None) or str(getattr(obj, 'pk', '')),
            ip_address=_client_ip(self.request),
            payload={
                'title': getattr(obj, 'title', None),
                'is_active': getattr(obj, 'is_active', None),
            },
        )
